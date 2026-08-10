#!/usr/bin/env python3
"""Generate deterministic browser preview fixtures from the bundled Office samples.

The application never executes LibreOffice at runtime. This developer-only helper
converts the checked-in samples to PDF and extracts a safe workbook/slide manifest.
"""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import tempfile
from datetime import date, datetime, time
from pathlib import Path
from typing import Any
from urllib.parse import quote

from openpyxl import load_workbook
from openpyxl.cell.cell import Cell
from pypdf import PdfReader
from pptx import Presentation


DOCUMENTS = {
    "pdf": "pdf测试.pdf",
    "word": "word测试.docx",
    "spreadsheet": "测试excel.xlsx",
    "presentation": "测试ppt.pptx",
}


def json_value(value: Any) -> Any:
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    return str(value)


def rgb_color(color: Any) -> str | None:
    if color is None or getattr(color, "type", None) != "rgb":
        return None
    value = getattr(color, "rgb", None)
    if not value:
        return None
    normalized = str(value)[-6:].upper()
    return f"#{normalized}" if len(normalized) == 6 else None


def cell_style(cell: Cell) -> dict[str, Any] | None:
    fill = rgb_color(cell.fill.fgColor) if cell.fill.fill_type == "solid" else None
    font_color = rgb_color(cell.font.color)
    border = any(
        getattr(side, "style", None)
        for side in (cell.border.left, cell.border.right, cell.border.top, cell.border.bottom)
    )
    style = {
        "bold": bool(cell.font.bold),
        "italic": bool(cell.font.italic),
        "fontColor": font_color,
        "fill": fill,
        "horizontal": cell.alignment.horizontal,
        "vertical": cell.alignment.vertical,
        "wrapText": bool(cell.alignment.wrap_text),
        "numberFormat": cell.number_format if cell.number_format != "General" else None,
        "border": border,
    }
    return style if any(value not in (None, False) for value in style.values()) else None


def workbook_manifest(source: Path) -> dict[str, Any]:
    formula_book = load_workbook(source, data_only=False, read_only=False)
    value_book = load_workbook(source, data_only=True, read_only=False)
    sheets: list[dict[str, Any]] = []
    for sheet in formula_book.worksheets:
        value_sheet = value_book[sheet.title]
        cells: list[dict[str, Any]] = []
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is None and not cell.has_style:
                    continue
                cached_value = value_sheet[cell.coordinate].value
                formula = cell.value if cell.data_type == "f" else None
                display_value = cached_value if formula is not None else cell.value
                record: dict[str, Any] = {
                    "address": cell.coordinate,
                    "row": cell.row,
                    "column": cell.column,
                    "value": json_value(display_value if display_value is not None else formula),
                }
                if formula is not None:
                    record["formula"] = str(formula)
                style = cell_style(cell)
                if style:
                    record["style"] = style
                cells.append(record)

        column_widths = {
            key: float(dimension.width)
            for key, dimension in sheet.column_dimensions.items()
            if dimension.width is not None
        }
        row_heights = {
            str(index): float(dimension.height)
            for index, dimension in sheet.row_dimensions.items()
            if dimension.height is not None
        }
        sheets.append(
            {
                "id": f"sheet-{len(sheets) + 1}",
                "name": sheet.title,
                "maxRow": sheet.max_row,
                "maxColumn": sheet.max_column,
                "cells": cells,
                "merges": [str(item) for item in sheet.merged_cells.ranges],
                "columnWidths": column_widths,
                "rowHeights": row_heights,
                "hasCharts": bool(sheet._charts),
            }
        )
    formula_book.close()
    value_book.close()
    return {"version": 1, "sheets": sheets}


def presentation_manifest(source: Path) -> list[dict[str, Any]]:
    presentation = Presentation(source)
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_blocks: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            normalized = " ".join(str(text).split())
            if normalized:
                text_blocks.append(normalized)
        slides.append(
            {
                "index": index,
                "title": text_blocks[0][:120] if text_blocks else f"幻灯片 {index}",
                "description": text_blocks[1][:180] if len(text_blocks) > 1 else None,
            }
        )
    return slides


def convert_to_pdf(soffice: Path, source: Path, work_dir: Path) -> Path:
    profile = work_dir / f"profile-{source.stem}"
    profile.mkdir(parents=True, exist_ok=True)
    profile_url = f"file://{quote(str(profile.resolve()))}"
    result = subprocess.run(
        [
            str(soffice),
            "--headless",
            f"-env:UserInstallation={profile_url}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(work_dir),
            str(source),
        ],
        check=False,
        capture_output=True,
        text=True,
        timeout=120,
    )
    output = work_dir / f"{source.stem}.pdf"
    if result.returncode != 0 or not output.exists() or output.stat().st_size == 0:
        raise RuntimeError(
            f"LibreOffice conversion failed for {source.name}: "
            f"{result.stderr.strip() or result.stdout.strip()}"
        )
    return output


def pdf_metadata(source: Path) -> dict[str, Any]:
    reader = PdfReader(source)
    return {
        "pageCount": len(reader.pages),
        "hasOutline": bool(reader.outline),
        "encrypted": bool(reader.is_encrypted),
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--documents",
        type=Path,
        default=Path(__file__).resolve().parents[1] / "src/assets/document",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path(__file__).resolve().parents[1] / "src/assets/document/previews",
    )
    parser.add_argument("--soffice", type=Path, default=Path(shutil.which("soffice") or "soffice"))
    args = parser.parse_args()

    documents = args.documents.resolve()
    output = args.output.resolve()
    output.mkdir(parents=True, exist_ok=True)
    sources = {key: documents / filename for key, filename in DOCUMENTS.items()}
    missing = [str(path) for path in sources.values() if not path.exists()]
    if missing:
        raise FileNotFoundError(f"Missing document fixtures: {missing}")

    with tempfile.TemporaryDirectory(prefix="aestival-document-previews-") as temp:
        work_dir = Path(temp)
        word_pdf = convert_to_pdf(args.soffice, sources["word"], work_dir)
        presentation_pdf = convert_to_pdf(args.soffice, sources["presentation"], work_dir)
        spreadsheet_pdf = convert_to_pdf(args.soffice, sources["spreadsheet"], work_dir)

        generated = {
            "word": output / "word测试.preview.pdf",
            "presentation": output / "测试ppt.preview.pdf",
            "spreadsheet": output / "测试excel.print.pdf",
        }
        shutil.copyfile(word_pdf, generated["word"])
        shutil.copyfile(presentation_pdf, generated["presentation"])
        shutil.copyfile(spreadsheet_pdf, generated["spreadsheet"])

    workbook = workbook_manifest(sources["spreadsheet"])
    slides = presentation_manifest(sources["presentation"])
    (output / "测试excel.workbook.json").write_text(
        json.dumps(workbook, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    (output / "测试ppt.slides.json").write_text(
        json.dumps({"version": 1, "slides": slides}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    manifest = {
        "version": 1,
        "documents": {
            DOCUMENTS["pdf"]: {"kind": "pdf", **pdf_metadata(sources["pdf"])},
            DOCUMENTS["word"]: {
                "kind": "word",
                "preview": generated["word"].name,
                **pdf_metadata(generated["word"]),
            },
            DOCUMENTS["presentation"]: {
                "kind": "presentation",
                "preview": generated["presentation"].name,
                "slideCount": len(slides),
                **pdf_metadata(generated["presentation"]),
            },
            DOCUMENTS["spreadsheet"]: {
                "kind": "spreadsheet",
                "preview": generated["spreadsheet"].name,
                "workbook": "测试excel.workbook.json",
                "sheetCount": len(workbook["sheets"]),
                "hasCharts": any(sheet["hasCharts"] for sheet in workbook["sheets"]),
                **pdf_metadata(generated["spreadsheet"]),
            },
        },
    }
    (output / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
